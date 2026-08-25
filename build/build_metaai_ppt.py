"""Build a light, elegant PowerPoint explaining the MetaAI paper.

Paper: "Enabling Over-the-Air AI for Edge Computing via Metasurface-Driven
Physical Neural Networks" (MetaAI), SIGCOMM '25, Feng et al.

The deck mirrors the companion Word document: it explains the idea for a
newcomer and adds the technical substance for a researcher, with the key
equations shown and their variables named. No emojis are used.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "assets_metaai")
OUT = os.path.join(os.path.dirname(HERE), "MetaAI_OverTheAir_Edge_AI.pptx")

# Palette - light and elegant
SLATE = RGBColor(0x4A, 0x6B, 0x8A)
SLATE_LT = RGBColor(0x7C, 0x97, 0xB2)
ACCENT = RGBColor(0x7A, 0x9A, 0x7E)
PLUM = RGBColor(0x7A, 0x64, 0x9A)
INK = RGBColor(0x3A, 0x3E, 0x44)
MUTE = RGBColor(0x8A, 0x8F, 0x96)
RULE = RGBColor(0xCF, 0xDA, 0xE4)
BG = RGBColor(0xFC, 0xFD, 0xFE)
HEADER_FILL = RGBColor(0xE7, 0xEE, 0xF4)
ROW_FILL = RGBColor(0xF5, 0xF8, 0xFA)
GREEN_FILL = RGBColor(0xED, 0xF4, 0xED)
PLUM_FILL = RGBColor(0xF1, 0xED, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set_font(run, size, color, bold=False, italic=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def rule(s, x, y, w, color=RULE, weight=1.5):
    ln = s.shapes.add_connector(2, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def accent_tab(s, x, y, h=Inches(0.42), w=Inches(0.07), color=SLATE):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def title(s, text, kicker=None):
    accent_tab(s, Inches(0.6), Inches(0.55))
    tb, tf = textbox(s, Inches(0.85), Inches(0.42), Inches(11.9), Inches(0.8))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    _set_font(r, 29, SLATE, bold=True)
    if kicker:
        p2 = tf.add_paragraph()
        rk = p2.add_run(); rk.text = kicker
        _set_font(rk, 13, SLATE_LT, italic=True)
    rule(s, Inches(0.62), Inches(1.42), Inches(12.1))


def footer(s, idx):
    tb, tf = textbox(s, Inches(0.62), Inches(7.02), Inches(11.0), Inches(0.32))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "MetaAI  \u00b7  Over-the-Air AI via Metasurface-Driven Physical Neural Networks"
    _set_font(r, 9, MUTE)
    tb2, tf2 = textbox(s, Inches(12.2), Inches(7.02), Inches(0.7), Inches(0.32))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(idx)
    _set_font(r2, 9, MUTE)


def bullets(s, items, x, y, w, h, size=17, gap=10):
    tb, tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        text = it["t"]
        strong = it.get("b", False)
        sub = it.get("sub", False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        if not sub:
            m = p.add_run(); m.text = "\u2014  "
            _set_font(m, size, ACCENT if not strong else SLATE, bold=False)
        else:
            sp = p.add_run(); sp.text = "      \u2022  "
            _set_font(sp, size - 2, SLATE_LT)
        r = p.add_run(); r.text = text
        _set_font(r, size if not sub else size - 2, INK if not strong else SLATE, bold=strong)
    return tb


def picture(s, img, x, y, w=None, h=None):
    path = os.path.join(ASSETS, img)
    with Image.open(path) as im:
        iw, ih = im.width, im.height
    ar = iw / ih
    if w is not None and h is None:
        return s.shapes.add_picture(path, x, y, width=w)
    if h is not None and w is None:
        return s.shapes.add_picture(path, x, y, height=h)
    return s.shapes.add_picture(path, x, y, width=w, height=h)


def picture_centered(s, img, cx, y, w):
    """Place a picture centred horizontally at cx with given width."""
    pic = picture(s, img, cx - w // 2, y, w=w)
    return pic


def caption(s, text, x, y, w):
    tb, tf = textbox(s, x, y, w, Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_font(r, 10, MUTE, italic=True)


def chip(s, x, y, w, h, label, body, fill, ec):
    """A soft rounded panel with a bold label and body text."""
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = ec; box.line.width = Pt(1.0)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    _set_font(r, 14, SLATE, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = body
    _set_font(r2, 12, INK)
    return box


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
band.fill.solid(); band.fill.fore_color.rgb = SLATE; band.line.fill.background(); band.shadow.inherit = False

tb, tf = textbox(s, Inches(0.9), Inches(2.05), Inches(11.8), Inches(2.2))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Computing in Thin Air"
_set_font(r, 46, SLATE, bold=True)
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "How MetaAI Turns the Wireless Channel into a Neural Network"
_set_font(r2, 21, SLATE_LT)
rule(s, Inches(0.95), Inches(4.25), Inches(8.2), color=RULE, weight=1.5)
tb3, tf3 = textbox(s, Inches(0.95), Inches(4.4), Inches(11.4), Inches(1.2))
p3 = tf3.paragraphs[0]
r3 = p3.add_run(); r3.text = "A guided explanation of the SIGCOMM 2025 paper by Feng, Liang, Li, Zhao, Jing, Xie and Chen"
_set_font(r3, 14, INK, italic=True)
p4 = tf3.add_paragraph()
r4 = p4.add_run(); r4.text = "Northwest University  \u00b7  University at Buffalo (SUNY)"
_set_font(r4, 13, INK, italic=True)
tb5, tf5 = textbox(s, Inches(0.95), Inches(6.3), Inches(11.0), Inches(0.5))
p5 = tf5.paragraphs[0]
r5 = p5.add_run(); r5.text = "Explained for the newcomer and for the researcher"
_set_font(r5, 12, MUTE)

# ===========================================================================
# 2. THE ONE-LINE IDEA
# ===========================================================================
s = slide()
title(s, "The Idea in One Line", "What MetaAI does")
bullets(s, [
    {"t": "Normally a small device sends raw data to a server, which then runs the AI \u2014 sending and computing are two separate costs.", "b": False},
    {"t": "MetaAI makes the wireless signal do the AI while it travels, so the server receives the answer, not the raw data.", "b": True},
    {"t": "A programmable metasurface in the room reshapes the passing signal exactly like a neural network's multiply-and-add.", "b": False},
    {"t": "Communication and computation become one single step.", "b": True},
], Inches(0.85), Inches(1.8), Inches(11.7), Inches(2.8), size=18, gap=13)
picture_centered(s, "fig_paradigms.png", SW // 2, Inches(4.4), Inches(9.0))
footer(s, 2)

# ===========================================================================
# 3. WHY IT MATTERS
# ===========================================================================
s = slide()
title(s, "Why It Matters", "The edge-AI trade-off MetaAI attacks")
chip(s, Inches(0.7), Inches(1.75), Inches(3.9), Inches(2.0), "On-device AI",
     "Fast answers, but small IoT devices lack the compute and battery to run it.", GREEN_FILL, ACCENT)
chip(s, Inches(4.75), Inches(1.75), Inches(3.9), Inches(2.0), "Transmit then compute",
     "The usual fix: send raw data to a server. Communication and computation stay separate costs.", ROW_FILL, SLATE_LT)
chip(s, Inches(8.8), Inches(1.75), Inches(3.85), Inches(2.0), "MetaAI",
     "Fuse the two: compute inside the wireless channel during transmission.", PLUM_FILL, PLUM)
bullets(s, [
    {"t": "Earlier physical neural networks compute at the speed of light, but the data must still be fully transmitted first \u2014 they are co-processors.", "b": False},
    {"t": "Over-the-air computing gave addition for free, but multiplication needed complex, expensive precoding unavailable to commodity IoT devices.", "b": False},
    {"t": "The gap: a simple, single-metasurface, end-to-end network that runs over an ordinary wireless link. MetaAI fills it.", "b": True},
], Inches(0.85), Inches(4.15), Inches(11.7), Inches(2.6), size=16.5, gap=12)
footer(s, 3)

# ===========================================================================
# 4. THREE FACTS JOINED
# ===========================================================================
s = slide()
title(s, "The Insight: Three Everyday Facts", "Join them and the channel becomes a neural network")
chip(s, Inches(0.7), Inches(1.85), Inches(3.85), Inches(2.3), "1. Radio is waves",
     "Signals travel as waves, like ripples on water.", GREEN_FILL, ACCENT)
chip(s, Inches(4.75), Inches(1.85), Inches(3.85), Inches(2.3), "2. The channel transforms them",
     "Walls and people scale, delay and echo the wave. That rule is 'the channel'.", ROW_FILL, SLATE_LT)
chip(s, Inches(8.8), Inches(1.85), Inches(3.85), Inches(2.3), "3. A network is multiply-add",
     "A linear layer just multiplies inputs by weights and sums them.", PLUM_FILL, PLUM)
bullets(s, [
    {"t": "'Multiply the signal by something and add up the results' is exactly what a channel already does.", "b": True},
    {"t": "A metasurface lets us choose that 'something' \u2014 so the channel carries out the network's weights.", "b": False},
], Inches(0.85), Inches(4.5), Inches(11.7), Inches(2.0), size=18, gap=13)
footer(s, 4)

# ===========================================================================
# 5. CHALLENGE 1: SEQUENTIAL VS PARALLEL
# ===========================================================================
s = slide()
title(s, "Challenge 1: Sequential vs Parallel", "The linearity trick")
bullets(s, [
    {"t": "Wireless sends data one symbol at a time; a network wants all inputs at once.", "b": False},
    {"t": "Because the network is linear, the parallel weighted sum can be built one product at a time.", "b": True},
    {"t": "Send x1 now, x2 next, keep a running total \u2014 the final total is identical to computing in parallel.", "b": False},
], Inches(0.85), Inches(1.75), Inches(5.3), Inches(4.6), size=17, gap=15)
picture(s, "fig_sequential.png", Inches(6.3), Inches(2.1), w=Inches(6.6))
caption(s, "A parallel weighted sum equals one product per time slot, accumulated.", Inches(6.3), Inches(5.5), Inches(6.6))
footer(s, 5)

# ===========================================================================
# 6. CHALLENGE 2: MULTIPLY IN THE WAVE
# ===========================================================================
s = slide()
title(s, "Challenge 2: Multiply in the Wave Itself", "The channel is a linear operator")
bullets(s, [
    {"t": "As a signal travels, the channel linearly transforms it:", "b": False},
], Inches(0.85), Inches(1.75), Inches(11.6), Inches(0.7), size=18)
picture_centered(s, "eq_channel.png", SW // 2, Inches(2.5), Inches(3.0))
bullets(s, [
    {"t": "x(t): the transmitted symbol (one input value).   H(t): the channel (the weight).   y(t): the product that arrives.", "b": False},
    {"t": "Program the metasurface so H(t) equals the weight the network wants at each instant.", "b": True},
    {"t": "Radio symbols are complex (amplitude and phase), so H(t) is a complex weight \u2014 hence a complex-valued network.", "b": False},
], Inches(0.85), Inches(3.9), Inches(11.6), Inches(2.6), size=16.5, gap=12)
footer(s, 6)

# ===========================================================================
# 7. CORE EQUATION
# ===========================================================================
s = slide()
title(s, "The Core MetaAI Equation", "Multiply over the air, accumulate in software")
picture_centered(s, "eq_core.png", SW // 2, Inches(1.75), Inches(4.6))
bullets(s, [
    {"t": "y_r: the score for class r.    H_r(t_i): the metasurface weight for class r at step i.    x_i: the i-th input symbol.", "b": False},
    {"t": "The sum is accumulated over time; the magnitude bars turn the complex result into a real score and give the nonlinearity.", "b": False},
    {"t": "Multiplication is physical (over the air); only the addition is digital, at the receiver.", "b": True},
    {"t": "Each input gets its own time slot, so one metasurface assigns an independent weight to every input \u2014 no stacked layers needed.", "b": False},
], Inches(0.85), Inches(3.35), Inches(11.7), Inches(3.2), size=16.5, gap=12)
footer(s, 7)

# ===========================================================================
# 8. HOW THE METASURFACE MAKES THE WEIGHT
# ===========================================================================
s = slide()
title(s, "How the Metasurface Makes a Weight", "From atom phases to channel value")
bullets(s, [
    {"t": "The channel through the metasurface is the sum of all meta-atom contributions:", "b": False},
], Inches(0.85), Inches(1.7), Inches(11.6), Inches(0.6), size=17)
picture_centered(s, "eq_mts.png", Inches(3.55), Inches(2.35), Inches(3.4))
picture(s, "fig_metasurface.png", Inches(7.2), Inches(2.2), w=Inches(5.7))
bullets(s, [
    {"t": "M: number of atoms (256).", "sub": True},
    {"t": "phi_m: the phase each atom adds \u2014 the only control knob.", "sub": True},
    {"t": "phi^p_m: the path phase.   alpha_p: a common amplitude that scales all outputs equally, so it does not change the class.", "sub": True},
], Inches(0.85), Inches(4.4), Inches(5.8), Inches(2.2), size=15, gap=9)
caption(s, "256 atoms, 4 phase states each (2-bit).", Inches(7.2), Inches(5.7), Inches(5.7))
footer(s, 8)

# ===========================================================================
# 9. SOLVING FOR THE SETTINGS
# ===========================================================================
s = slide()
title(s, "Solving for the Metasurface Settings", "Match the trained weight as closely as hardware allows")
bullets(s, [
    {"t": "Goal: the metasurface channel should equal the desired trained weight. Hardware cannot be exact, so minimise the gap:", "b": False},
], Inches(0.85), Inches(1.75), Inches(11.6), Inches(0.9), size=17)
picture_centered(s, "eq_config.png", SW // 2, Inches(2.75), Inches(7.6))
bullets(s, [
    {"t": "Phi: the atom phases we solve for.   H_mts: what they produce.   H_des: the target weight from training.", "b": False},
    {"t": "Each atom offers only 4 discrete phases; the optimiser searches them across all atoms.", "b": False},
    {"t": "Train continuous, then snap to discrete \u2014 this beats training with discrete weights from the start.", "b": True},
], Inches(0.85), Inches(4.25), Inches(11.7), Inches(2.4), size=16.5, gap=11)
footer(s, 9)

# ===========================================================================
# 10. WORKFLOW
# ===========================================================================
s = slide()
title(s, "End-to-End Workflow", "From trained model to computing in the air")
picture_centered(s, "fig_workflow.png", SW // 2, Inches(2.0), Inches(9.6))
bullets(s, [
    {"t": "Train a complex-valued linear network, read out its desired weights, solve for the atom phases, load them, and let the signal compute at run time.", "b": True},
    {"t": "The whole neural network lives in a single reconfigurable metasurface.", "b": False},
], Inches(0.85), Inches(4.7), Inches(11.7), Inches(2.0), size=17, gap=13)
footer(s, 10)

# ===========================================================================
# 11. REAL-ROOM ROBUSTNESS
# ===========================================================================
s = slide()
title(s, "Making It Work in a Real Room", "Four practical problems, four remedies")
chip(s, Inches(0.7), Inches(1.8), Inches(5.85), Inches(1.75), "Multipath echoes",
     "Symbols average to zero, so echoes cancel; the metasurface varies its weight mid-symbol so its signal survives.", GREEN_FILL, ACCENT)
chip(s, Inches(6.75), Inches(1.8), Inches(5.85), Inches(1.75), "Clock drift (CDFA)",
     "Coarse energy detection triggers the weights; training injects random shifts so the model tolerates the residual error.", ROW_FILL, SLATE_LT)
chip(s, Inches(0.7), Inches(3.75), Inches(5.85), Inches(1.75), "System noise",
     "Model hardware and environmental noise as a pre-disturbed signal; train at lower SNR so it copes at run time.", PLUM_FILL, PLUM)
chip(s, Inches(6.75), Inches(3.75), Inches(5.85), Inches(1.75), "Latency (parallelism)",
     "Subcarriers or antennas compute several classes at once, solved by one joint optimisation \u2014 speed for a small accuracy cost.", ROW_FILL, SLATE_LT)
bullets(s, [
    {"t": "Each remedy was verified separately: CDFA lifts accuracy from ~19% to ~89%; multipath cancellation holds >82% across three rooms.", "b": True},
], Inches(0.7), Inches(5.75), Inches(12.0), Inches(1.0), size=14.5, gap=6)
footer(s, 11)

# ===========================================================================
# 12. MULTIPATH FIGURE (deep dive)
# ===========================================================================
s = slide()
title(s, "Multipath Cancellation, Visualised", "Break the zero-mean property for the wanted path only")
picture_centered(s, "fig_multipath.png", SW // 2, Inches(1.85), Inches(10.5))
bullets(s, [
    {"t": "(a) A symbol averages to zero over its period.  (b) Environmental echoes keep that zero average, so they cancel when summed.", "b": False},
    {"t": "(c) The metasurface changes its weight within the symbol, so its contribution does not cancel \u2014 it is preserved. No channel estimation needed.", "b": True},
], Inches(0.85), Inches(4.75), Inches(11.7), Inches(2.0), size=16, gap=12)
footer(s, 12)

# ===========================================================================
# 13. MULTI-SENSOR FUSION
# ===========================================================================
s = slide()
title(s, "Multi-Sensor Fusion", "One shared metasurface, many sensors")
bullets(s, [
    {"t": "Weights for different sensors are independent, so each sensor is processed in its own time window on the same surface.", "b": False},
    {"t": "Per-sensor scores are summed (late fusion): complementary sensors reinforce the right answer.", "b": True},
    {"t": "No extra hardware per sensor \u2014 unlike traditional PNNs.", "b": False},
], Inches(0.85), Inches(1.75), Inches(5.5), Inches(3.4), size=16.5, gap=13)
picture_centered(s, "eq_fusion.png", Inches(3.5), Inches(5.2), Inches(3.2))
picture(s, "fig_multisensor.png", Inches(6.5), Inches(2.0), w=Inches(6.3))
caption(s, "Fusing views, antennas or modalities raises accuracy by up to 27%.", Inches(6.5), Inches(5.7), Inches(6.3))
footer(s, 13)

# ===========================================================================
# 14. HARDWARE
# ===========================================================================
s = slide()
title(s, "The Hardware", "A built and measured prototype")
rows_data = [
    ["Component", "What the authors used"],
    ["Metasurfaces", "Dual-band (2.4 / 5 GHz) and single-band (3.5 GHz)"],
    ["Meta-atom grid", "16 x 16 = 256 atoms per surface"],
    ["Per-atom control", "Two PIN diodes: 4 phase states (2-bit)"],
    ["Radios", "USRP X310 software-defined radios (Tx and Rx)"],
    ["Default link", "5.25 GHz, 256-QAM, 1 M symbols/s"],
    ["Environments", "Corridor, lab, office; NLoS and cross-room"],
]
gtbl = s.shapes.add_table(len(rows_data), 2, Inches(0.7), Inches(1.75), Inches(6.4), Inches(4.4)).table
gtbl.columns[0].width = Inches(2.3); gtbl.columns[1].width = Inches(4.1)
for ri, row in enumerate(rows_data):
    gtbl.rows[ri].height = Inches(0.6)
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL if ri == 0 else (ROW_FILL if ri % 2 == 0 else WHITE)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val
        _set_font(r, 12.5 if ri == 0 else 12, INK if ri else SLATE, bold=(ri == 0 or ci == 0))
picture(s, "fig_metaatoms.png", Inches(7.35), Inches(2.15), w=Inches(5.5))
caption(s, "Accuracy and weight density both saturate near 256 atoms \u2014 the chosen size.",
        Inches(7.35), Inches(5.75), Inches(5.5))
footer(s, 14)

# ===========================================================================
# 15. RESULTS
# ===========================================================================
s = slide()
title(s, "What the Experiments Show", "Strong accuracy from a single linear layer")
picture(s, "fig_accuracy.png", Inches(0.7), Inches(1.75), w=Inches(6.6))
caption(s, "MetaAI stays close to a digital ResNet-18 and beats a discrete-from-start baseline.",
        Inches(0.7), Inches(5.75), Inches(6.6))
bullets(s, [
    {"t": "82.8% average accuracy across six datasets, up to 89.8% on digits.", "b": True},
    {"t": "Within ~7% of its own digital simulation \u2014 a faithful physical build.", "b": False},
    {"t": "Beats the DiscreteNN baseline on every dataset.", "b": False},
    {"t": "Multi-sensor fusion adds up to 27%.", "b": True},
    {"t": "Robust across bands (2.4/3.5/5 GHz), modulations, distances, angles, NLoS and cross-room.", "b": False},
], Inches(7.6), Inches(1.9), Inches(5.1), Inches(4.6), size=15.5, gap=12)
footer(s, 15)

# ===========================================================================
# 16. STRENGTHS & LIMITATIONS
# ===========================================================================
s = slide()
title(s, "Strengths, Limitations, Takeaway", "An honest balance")
bullets(s, [
    {"t": "Strengths", "b": True},
    {"t": "Saves IoT energy and cost \u2014 the device only transmits.", "sub": True},
    {"t": "Private by design \u2014 the server sees results, not raw data.", "sub": True},
    {"t": "One shared metasurface; works over standard commodity links.", "sub": True},
], Inches(0.85), Inches(1.75), Inches(5.9), Inches(4.6), size=16, gap=10)
bullets(s, [
    {"t": "Limitations", "b": True},
    {"t": "Linear networks only; nonlinear/deeper models are future work.", "sub": True},
    {"t": "Sequential input ties model size to the latency budget.", "sub": True},
    {"t": "Precision capped by atom count and 2-bit depth.", "sub": True},
    {"t": "Motion needs recalibration of the phase-to-weight map.", "sub": True},
], Inches(7.0), Inches(1.75), Inches(5.9), Inches(4.6), size=16, gap=10)
chip(s, Inches(0.85), Inches(5.55), Inches(12.05), Inches(1.15), "Takeaway",
     "MetaAI turns the wireless channel from a passive pipe into an active computer, unifying communication and computation "
     "for always-on, low-power, privacy-sensitive edge sensing.", PLUM_FILL, PLUM)
footer(s, 16)

# ===========================================================================
# 17. CLOSING
# ===========================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
band.fill.solid(); band.fill.fore_color.rgb = SLATE; band.line.fill.background(); band.shadow.inherit = False
tb, tf = textbox(s, Inches(0.95), Inches(2.7), Inches(11.6), Inches(2.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Compute where the signal already goes."
_set_font(r, 34, SLATE, bold=True)
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "MetaAI: a single metasurface that is, at once, the network's weights and the wireless link."
_set_font(r2, 17, SLATE_LT, italic=True)
rule(s, Inches(1.0), Inches(4.7), Inches(7.4))
tb3, tf3 = textbox(s, Inches(1.0), Inches(4.85), Inches(11.0), Inches(0.6))
p3 = tf3.paragraphs[0]
r3 = p3.add_run(); r3.text = "Companion deck to the SIGCOMM 2025 paper by Feng et al."
_set_font(r3, 13, MUTE)
footer(s, 17)

prs.save(OUT)
print("wrote", OUT)
