# -*- coding: utf-8 -*-
"""
7-slide deck for the 22 July meeting, covering BOTH research tracks:
  1. ESP32 WiFi (CSI) sensing hardware trial  (project already run)
  2. MetaAI over-the-air compute reproduction  (Direction J from the dossier)

Content for track 2 is taken from WirelessPerception_Dossier 3.docx:
  - red-highlighted direction  = Direction J (over-the-air compute on Wi-Fi CSI)
  - yellow-highlighted papers   = the ones I can use / liked
  - J1 "Proposed project plan"  = my proposed plan + milestones
Clean, plain language, nothing invented.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

MEET = r"c:\Users\fs161326\Downloads\SNU Research\22 july meet"

IMG = {
    "grid":       os.path.join(MEET, "WhatsApp Image 2026-07-22 at 01.42.21.jpeg"),
    "timeseries": os.path.join(MEET, "WhatsApp Image 2026-07-22 at 01.10.43.jpeg"),
    "fingerprint":os.path.join(MEET, "WhatsApp Image 2026-07-22 at 01.10.433.jpeg"),
}

INK   = RGBColor(0x1F, 0x2A, 0x37)
BLUE  = RGBColor(0x2E, 0x5A, 0x88)
SLATE = RGBColor(0x54, 0x63, 0x72)
ACCENT= RGBColor(0xC0, 0x4A, 0x3A)
GREEN = RGBColor(0x2E, 0x6B, 0x3A)
BG    = RGBColor(0xF7, 0xF7, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def line(tf, text, size=18, color=INK, bold=False, italic=False,
         space_after=8, bullet=False, align=PP_ALIGN.LEFT, first=False, indent=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        text = ("      \u2013  " if indent else "\u2022  ") + text
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"
    return p


def sidebar(slide, color=BLUE):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def header(slide, kicker, title, kcolor=ACCENT, bcolor=BLUE):
    sidebar(slide, bcolor)
    tf = tb(slide, Inches(0.6), Inches(0.42), Inches(12.2), Inches(1.25))
    line(tf, kicker.upper(), size=13, color=kcolor, bold=True, space_after=2, first=True)
    line(tf, title, size=29, color=INK, bold=True, space_after=0)


def caption(slide, l, t, w, text):
    tf = tb(slide, l, t, w, Inches(0.4))
    line(tf, text, size=12, color=SLATE, italic=True, space_after=0, first=True,
         align=PP_ALIGN.CENTER)


def add_img_fit(slide, path, l, t, max_w, max_h):
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1600, 900
    ar = iw / ih; box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    slide.shapes.add_picture(path, l + (max_w - w) // 2, t + (max_h - h) // 2, width=w, height=h)


# ============================================================ Slide 1: title
s = prs.slides.add_slide(BLANK); bg(s, INK)
band = s.shapes.add_shape(1, 0, Inches(2.55), SW, Inches(0.06))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
tf = tb(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.4))
line(tf, "Wireless perception: two tracks",
     size=42, color=WHITE, bold=True, space_after=10, first=True)
line(tf, "1 \u00b7 A hands-on ESP32 WiFi sensing trial      2 \u00b7 Reproducing the MetaAI \u201ccompute-in-the-channel\u201d idea for WiFi",
     size=19, color=RGBColor(0xC9, 0xD3, 0xDE), space_after=0)
tf = tb(s, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.8))
line(tf, "MPL Research  \u00b7  22 July meeting",
     size=14, color=RGBColor(0x8F, 0x9C, 0xAB), space_after=0, first=True)

# ============================================================ Slide 2: two directions
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "The two briefs from the lab", "What I looked at, and what I am proposing")

c1 = s.shapes.add_shape(1, Inches(0.6), Inches(1.95), Inches(5.9), Inches(4.7))
c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
c1.line.color.rgb = RGBColor(0xD8, 0xDD, 0xE2); c1.line.width = Pt(1); c1.shadow.inherit = False
tf = c1.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.28)
line(tf, "TRACK 1  \u00b7  HARDWARE", size=13, color=ACCENT, bold=True, space_after=6, first=True)
line(tf, "WiFi sensing on real ESP32 boards", size=19, color=INK, bold=True, space_after=12)
line(tf, "Can cheap WiFi boards sense a person\u2019s motion and rough location, "
         "using camera labels for ground truth?", size=15, space_after=10)
line(tf, "Status: built and run \u2014 first results in this deck.", size=15, color=GREEN, bold=True, space_after=0)

c2 = s.shapes.add_shape(1, Inches(6.8), Inches(1.95), Inches(5.9), Inches(4.7))
c2.fill.solid(); c2.fill.fore_color.rgb = WHITE
c2.line.color.rgb = RGBColor(0xD8, 0xDD, 0xE2); c2.line.width = Pt(1); c2.shadow.inherit = False
tf = c2.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.28)
line(tf, "TRACK 2  \u00b7  THE METAAI PAPER", size=13, color=ACCENT, bold=True, space_after=6, first=True)
line(tf, "\u201cComputing in the channel\u201d for WiFi", size=19, color=INK, bold=True, space_after=12)
line(tf, "MetaAI does neural-network computation in the wireless channel itself. "
         "The brief: implement it for wireless perception tasks.", size=15, space_after=10)
line(tf, "Status: read the literature, chosen a direction, drafted a plan.", size=15, color=BLUE, bold=True, space_after=0)

# ============================================================ Slide 3: ESP32 setup+method
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Track 1 \u00b7 the hardware trial", "Two ESP32 boards + a camera for ground truth")

tf = tb(s, Inches(0.6), Inches(1.95), Inches(6.3), Inches(4.8))
line(tf, "The setup", size=18, color=BLUE, bold=True, space_after=8, first=True)
line(tf, "Two ESP32-WROOM-32 boards (ESPectre firmware). Each reads the WiFi channel "
         "and outputs a Movement Score, 0\u201310.", size=15, bullet=True, space_after=6)
line(tf, "One board keeps WiFi traffic flowing (~100 pings/sec) so there is always "
         "a signal to measure.", size=15, bullet=True, space_after=6)
line(tf, "A camera detects the person and labels which zone they stand in. Room split "
         "into a 4 \u00d7 3 grid.", size=15, bullet=True, space_after=14)
line(tf, "The method", size=18, color=BLUE, bold=True, space_after=8)
line(tf, "When a person moves, the score rises; when still, it settles.", size=15, bullet=True, space_after=6)
line(tf, "Zone label + both scores + position saved together in one CSV.", size=15, bullet=True, space_after=6)
line(tf, "This run: ~4,660 samples over ~3.5 minutes, labelled automatically.", size=15, bullet=True, space_after=0)

add_img_fit(s, IMG["grid"], Inches(7.35), Inches(2.0), Inches(5.45), Inches(4.4))
caption(s, Inches(7.35), Inches(6.5), Inches(5.45),
        "Live camera view with the 4\u00d73 zone grid overlaid")

# ============================================================ Slide 4: ESP32 results + honest read
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Track 1 \u00b7 results", "It reacts to motion \u2014 honestly, not solved yet")

add_img_fit(s, IMG["fingerprint"], Inches(0.5), Inches(1.9), Inches(6.9), Inches(4.5))
caption(s, Inches(0.5), Inches(6.45), Inches(6.9),
        "Average score per zone \u2014 ESP32-1 pinned near 10, ESP32-2 varies by zone")

tf = tb(s, Inches(7.7), Inches(2.0), Inches(5.1), Inches(4.6))
line(tf, "What works", size=17, color=GREEN, bold=True, space_after=6, first=True)
line(tf, "Full pipeline runs end to end: sense \u2192 camera-label \u2192 save.", size=14.5, bullet=True, space_after=5)
line(tf, "Both boards genuinely pick up human motion.", size=14.5, bullet=True, space_after=5)
line(tf, "ESP32-2 gives a graded, position-dependent response.", size=14.5, bullet=True, space_after=12)
line(tf, "What needs fixing", size=17, color=ACCENT, bold=True, space_after=6)
line(tf, "ESP32-1 saturates near 10 \u2014 can\u2019t tell zones apart yet; needs threshold / "
         "placement / calibration tuning.", size=14.5, bullet=True, space_after=5)
line(tf, "Coverage was uneven (mostly one row), so per-zone averages aren\u2019t fair yet.",
     size=14.5, bullet=True, space_after=5)
line(tf, "So far a solid proof of concept, not a trained location model.", size=14.5, bullet=True, space_after=0)

# ============================================================ Slide 5: MetaAI - what I read
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Track 2 \u00b7 what I read", "The direction I picked, and the papers I liked",
       bcolor=ACCENT)

tf = tb(s, Inches(0.6), Inches(1.95), Inches(6.4), Inches(4.8))
line(tf, "Chosen direction", size=18, color=ACCENT, bold=True, space_after=6, first=True)
line(tf, "Over-the-air compute on WiFi CSI", size=17, color=INK, bold=True, space_after=6)
line(tf, "Take MetaAI\u2019s software recipe (linear decomposition, CDFA sync, noise-aware "
         "training, multi-sensor fusion) and test it \u2014 in simulation \u2014 on real "
         "through-wall WiFi datasets (Widar3.0, Person-in-WiFi 3D).", size=14.5, space_after=6)
line(tf, "Nobody has done this comparison yet, and it needs no hardware \u2014 which is "
         "why I picked it.", size=14.5, italic=True, color=SLATE, space_after=0)

tf = tb(s, Inches(7.3), Inches(1.95), Inches(5.5), Inches(4.8))
line(tf, "Papers I can use / liked", size=18, color=ACCENT, bold=True, space_after=8, first=True)
line(tf, "MINN \u2014 over-the-air edge inference via metasurface-integrated neural nets.", size=14.5, bullet=True, space_after=6)
line(tf, "AirCNN \u2014 CNNs computed in the air via reconfigurable surfaces.", size=14.5, bullet=True, space_after=6)
line(tf, "RISAR \u2014 clearest proof a controllable surface boosts WiFi sensing (a baseline to beat).", size=14.5, bullet=True, space_after=6)
line(tf, "LWM \u2014 a foundation model for wireless channels.", size=14.5, bullet=True, space_after=6)
line(tf, "WhoFi \u2014 person re-identification from WiFi channel signals.", size=14.5, bullet=True, space_after=0)

# ============================================================ Slide 6: MetaAI - proposed plan
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Track 2 \u00b7 my proposed plan", "A software reproduction, benchmarked honestly",
       bcolor=ACCENT)

tf = tb(s, Inches(0.6), Inches(1.95), Inches(6.5), Inches(4.8))
line(tf, "What I would build", size=18, color=ACCENT, bold=True, space_after=8, first=True)
line(tf, "A Python simulator: 256-cell 2-bit metasurface + a time-varying channel "
         "(MetaAI\u2019s equations).", size=14.5, bullet=True, space_after=6)
line(tf, "A trainable complex-valued linear network with MetaAI\u2019s sync (CDFA) and "
         "noise augmentations.", size=14.5, bullet=True, space_after=6)
line(tf, "A benchmark on Widar3.0 & Person-in-WiFi 3D against a standard CSI classifier, "
         "Wi-CBR, UniCrossFi and an RF-Diffusion baseline.", size=14.5, bullet=True, space_after=6)
line(tf, "An honest report \u2014 including negative results \u2014 on where MetaAI\u2019s tricks "
         "help and where they hurt.", size=14.5, bullet=True, space_after=0)

mc = s.shapes.add_shape(1, Inches(7.4), Inches(2.0), Inches(5.4), Inches(4.5))
mc.fill.solid(); mc.fill.fore_color.rgb = WHITE
mc.line.color.rgb = RGBColor(0xD8, 0xDD, 0xE2); mc.line.width = Pt(1); mc.shadow.inherit = False
tf = mc.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.32); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.28)
line(tf, "Rough milestones", size=17, color=BLUE, bold=True, space_after=12, first=True)
for wk, txt in [
    ("Wk 1\u20132", "reproduce MetaAI\u2019s MNIST baseline in the simulator"),
    ("Wk 3\u20134", "port to Widar3.0, get a first-pass number"),
    ("Wk 5\u20136", "add CDFA sync, multipath, noise-aware training"),
    ("Wk 7\u201310", "add Wi-CBR / UniCrossFi / RF-Diffusion baselines + ablations"),
    ("Wk 11\u201314", "write up as a workshop paper"),
]:
    p = tf.add_paragraph(); p.space_after = Pt(9)
    r = p.add_run(); r.text = wk + "  "
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(14); r2.font.color.rgb = INK; r2.font.name = "Calibri"

# ============================================================ Slide 7: next steps both
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Where both tracks go", "Next steps")

tf = tb(s, Inches(0.6), Inches(1.95), Inches(5.9), Inches(4.8))
line(tf, "Track 1 \u2014 ESP32", size=18, color=BLUE, bold=True, space_after=8, first=True)
line(tf, "Re-tune ESP32-1 so both boards give comparable graded responses.", size=15, bullet=True, space_after=7)
line(tf, "Guided walk with equal time in every zone for fair fingerprints.", size=15, bullet=True, space_after=7)
line(tf, "Train a simple zone classifier from the two scores + camera labels.", size=15, bullet=True, space_after=0)

tf = tb(s, Inches(6.9), Inches(1.95), Inches(5.9), Inches(4.8))
line(tf, "Track 2 \u2014 MetaAI", size=18, color=ACCENT, bold=True, space_after=8, first=True)
line(tf, "Stand up the simulator and reproduce the MNIST baseline first.", size=15, bullet=True, space_after=7)
line(tf, "Move to Widar3.0 and get an honest first number.", size=15, bullet=True, space_after=7)
line(tf, "Both tracks share one goal: perception from the wireless signal, "
         "reducing reliance on cameras.", size=15, bullet=True, space_after=0)

out = os.path.join(MEET, "ESP32_WiFi_Sensing_22July.pptx")
prs.save(out)
print("wrote", out)
