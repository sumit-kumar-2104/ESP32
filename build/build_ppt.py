"""Build an elegant, light PowerPoint deck mirroring the wireless wave-perception white paper."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(os.path.dirname(HERE), "Wireless_Wave_Perception_and_MetaPerception.pptx")

# Palette - light and elegant
SLATE = RGBColor(0x4A, 0x6B, 0x8A)
SLATE_LT = RGBColor(0x7C, 0x97, 0xB2)
ACCENT = RGBColor(0x7A, 0x9A, 0x7E)
INK = RGBColor(0x3A, 0x3E, 0x44)
MUTE = RGBColor(0x8A, 0x8F, 0x96)
RULE = RGBColor(0xCF, 0xDA, 0xE4)
BG = RGBColor(0xFC, 0xFD, 0xFE)
HEADER_FILL = RGBColor(0xE7, 0xEE, 0xF4)
ROW_FILL = RGBColor(0xF5, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set_font(run, size, color, bold=False, italic=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def rule(s, x, y, w, color=RULE, weight=1.5):
    ln = s.shapes.add_connector(2, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def accent_tab(s, x, y, h=Inches(0.42), w=Inches(0.07), color=SLATE):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def title(s, text, kicker=None):
    accent_tab(s, Inches(0.6), Inches(0.55))
    tb, tf = textbox(s, Inches(0.85), Inches(0.42), Inches(11.8), Inches(0.8))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    _set_font(r, 30, SLATE, bold=True)
    if kicker:
        p2 = tf.add_paragraph()
        rk = p2.add_run(); rk.text = kicker
        _set_font(rk, 13, SLATE_LT, italic=True)
    rule(s, Inches(0.62), Inches(1.42), Inches(12.1))


def footer(s, idx, label):
    tb, tf = textbox(s, Inches(0.62), Inches(7.02), Inches(11.0), Inches(0.32))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    _set_font(r, 9, MUTE)
    tb2, tf2 = textbox(s, Inches(12.2), Inches(7.02), Inches(0.7), Inches(0.32))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(idx)
    _set_font(r2, 9, MUTE)


def bullets(s, items, x, y, w, h, size=17, gap=10):
    tb, tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        lead = it.get("lead", 0)
        text = it["t"]
        strong = it.get("b", False)
        sub = it.get("sub", False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        # accent dash marker
        if not sub:
            m = p.add_run(); m.text = "\u2014  "
            _set_font(m, size, ACCENT if not strong else SLATE, bold=False)
        else:
            sp = p.add_run(); sp.text = "      \u2022  "
            _set_font(sp, size - 2, SLATE_LT)
        r = p.add_run(); r.text = text
        _set_font(r, size if not sub else size - 2, INK if not strong else SLATE, bold=strong)
    return tb


def picture(s, img, x, y, w=None, h=None):
    path = os.path.join(ASSETS, img)
    with Image.open(path) as im:
        iw, ih = im.width, im.height
    ar = iw / ih
    if w is not None and h is None:
        h = Emu(int(w * (1 / ar)))
        return s.shapes.add_picture(path, x, y, width=w)
    if h is not None and w is None:
        return s.shapes.add_picture(path, x, y, height=h)
    return s.shapes.add_picture(path, x, y, width=w, height=h)


def caption(s, text, x, y, w):
    tb, tf = textbox(s, x, y, w, Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_font(r, 10, MUTE, italic=True)


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
# soft accent band on the left
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SH)
band.fill.solid(); band.fill.fore_color.rgb = SLATE; band.line.fill.background(); band.shadow.inherit = False

tb, tf = textbox(s, Inches(0.9), Inches(2.25), Inches(11.6), Inches(2.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "From Light to Radio"
_set_font(r, 46, SLATE, bold=True)
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "A Mathematical Case for Wireless Wave-Based Perception"
_set_font(r2, 22, SLATE_LT)
rule(s, Inches(0.95), Inches(4.35), Inches(7.6), color=RULE, weight=1.5)
tb3, tf3 = textbox(s, Inches(0.95), Inches(4.5), Inches(11.0), Inches(1.0))
p3 = tf3.paragraphs[0]
r3 = p3.add_run(); r3.text = "Extending optical wave-computation and diffractive models to wireless sensing,"
_set_font(r3, 14, INK, italic=True)
p4 = tf3.add_paragraph()
r4 = p4.add_run(); r4.text = "with applications to multi-modal meta-perception"
_set_font(r4, 14, INK, italic=True)
tb5, tf5 = textbox(s, Inches(0.95), Inches(6.4), Inches(11.0), Inches(0.5))
p5 = tf5.paragraphs[0]
r5 = p5.add_run(); r5.text = "Multi-Modal Perception programme   \u00b7   April 2026"
_set_font(r5, 12, MUTE)

# ===========================================================================
# 2. THE QUESTION & CLAIM
# ===========================================================================
s = slide()
title(s, "The Question, and the Claim", "Can wireless models do for perception what optical models do?")
bullets(s, [
    {"t": "Optical models clearly work \u2014 diffractive networks and optical generative models perform perception and synthesis in the wave itself.", "b": False},
    {"t": "Claim A  (sensing):", "b": True},
    {"t": "Treat the measured radio channel as the \u201cimage\u201d; the same deep-learning machinery gives detection, pose, activity and identity.", "sub": True},
    {"t": "Claim B  (computation):", "b": True},
    {"t": "Programmable metasurfaces let the radio wave itself carry out the computation \u2014 the true analogue of the optical model.", "sub": True},
    {"t": "Both are supported by published, peer-reviewed work.", "b": False},
], Inches(0.85), Inches(1.85), Inches(11.6), Inches(4.8), size=18, gap=12)
footer(s, 2, "From Light to Radio  \u00b7  Wireless Wave-Based Perception")

# ===========================================================================
# 3. INTUITION: WAVES ARE WAVES
# ===========================================================================
s = slide()
title(s, "Waves Are Waves", "The intuition in one picture")
bullets(s, [
    {"t": "Ripples on a pond spread, bounce and interfere. Shaped obstacles sculpt the far-side pattern \u2014 that is how a wave computes.", "b": False},
    {"t": "Light does exactly this in a diffractive network; the answer forms physically as light crosses patterned surfaces.", "b": False},
    {"t": "Radio waves ripple and interfere the same way \u2014 and a person in a room already reshapes them.", "b": True},
], Inches(0.85), Inches(1.75), Inches(5.5), Inches(4.6), size=17, gap=14)
picture(s, "fig_spectrum.png", Inches(6.6), Inches(2.2), w=Inches(6.3))
caption(s, "One spectrum, one wave equation \u2014 only the wavelength differs.", Inches(6.6), Inches(5.25), Inches(6.3))
footer(s, 3, "From Light to Radio  \u00b7  Plain-language intuition")

# ===========================================================================
# 4. ONE WAVE EQUATION
# ===========================================================================
s = slide()
title(s, "One Equation Shared by Optics and Wireless", "The foundation")
bullets(s, [
    {"t": "From Maxwell\u2019s equations, every field component obeys the time-harmonic Helmholtz equation:", "b": False},
], Inches(0.85), Inches(1.8), Inches(11.6), Inches(0.9), size=18)
picture(s, "eq_helmholtz.png", Inches(3.4), Inches(2.95), w=Inches(6.4))
bullets(s, [
    {"t": "The equation knows nothing of \u201coptical\u201d or \u201cradio\u201d.", "b": True},
    {"t": "A visible photon (\u03bb \u2248 0.5 \u00b5m) and a Wi-Fi carrier (\u03bb \u2248 6 cm) differ only in the value of \u03bb.", "b": False},
    {"t": "Everything that follows is a consequence of this single fact.", "b": False},
], Inches(0.85), Inches(4.4), Inches(11.6), Inches(2.2), size=17, gap=10)
footer(s, 4, "From Light to Radio  \u00b7  Mathematical foundation")

# ===========================================================================
# 5. HOW A WAVE COMPUTES
# ===========================================================================
s = slide()
title(s, "How a Wave Computes", "Diffraction as a trainable network")
bullets(s, [
    {"t": "Propagation is a dense linear operator \u2014 every input point reaches every output point, like a fully-connected layer.", "b": False},
    {"t": "Thin trainable surfaces apply a phase profile; intensity at the detector supplies the nonlinearity.", "b": False},
    {"t": "The phase profiles are the weights, trained by gradient descent \u2014 identical at any wavelength.", "b": True},
], Inches(0.85), Inches(1.75), Inches(5.4), Inches(4.6), size=17, gap=14)
picture(s, "fig_diffractive.png", Inches(6.45), Inches(2.0), w=Inches(6.5))
caption(s, "Same trainable network, optical and radio \u2014 only the scale changes.", Inches(6.45), Inches(5.35), Inches(6.5))
footer(s, 5, "From Light to Radio  \u00b7  Diffractive networks")

# ===========================================================================
# 6. THE PROOF: SCALE INVARIANCE
# ===========================================================================
s = slide()
title(s, "The Proof: Scale Invariance", "Why wireless inherits all of it")
bullets(s, [
    {"t": "Scale all distances and the wavelength together by  s = \u03bb_RF / \u03bb_opt :", "b": False},
], Inches(0.85), Inches(1.8), Inches(11.6), Inches(0.8), size=18)
picture(s, "eq_scale.png", Inches(3.2), Inches(2.7), w=Inches(6.8))
bullets(s, [
    {"t": "The product k\u00b7r that controls all interference is unchanged, so propagation between scaled surfaces is identical.", "b": False},
    {"t": "Copying the same surfaces, an L-layer radio network computes the same function as the optical one.", "b": True},
    {"t": "This is the classical scale invariance of Maxwell\u2019s equations; metasurfaces synthesise the required surfaces directly.", "b": False},
], Inches(0.85), Inches(4.15), Inches(11.6), Inches(2.4), size=17, gap=11)
footer(s, 6, "From Light to Radio  \u00b7  Electromagnetic similitude")

# ===========================================================================
# 7. WHAT CHANGES: RESOLUTION
# ===========================================================================
s = slide()
title(s, "What Actually Changes: Resolution", "The one principled trade-off")
bullets(s, [
    {"t": "Scale invariance preserves the computation \u2014 but not the fineness of detail:", "b": False},
    {"t": "Radio \u03bb is 10\u2074\u201310\u2075\u00d7 larger, so features are that much coarser at equal aperture.", "b": True},
    {"t": "Result: presence, motion, pose, activity \u2014 not faces or fine texture.", "b": False},
    {"t": "Coarse, non-identifying sensing is exactly what private, always-on monitoring wants.", "b": False},
], Inches(0.85), Inches(1.75), Inches(5.6), Inches(4.6), size=17, gap=13)
picture(s, "fig_resolution.png", Inches(6.8), Inches(2.05), w=Inches(5.9))
caption(s, "Resolution scales with wavelength: \u03b4 \u2248 \u03bb / (2\u00b7NA).", Inches(6.8), Inches(5.65), Inches(5.9))
footer(s, 7, "From Light to Radio  \u00b7  The resolution trade-off")

# ===========================================================================
# 8. WIRELESS AS A MEASURED WAVE FIELD
# ===========================================================================
s = slide()
title(s, "Wireless Perception as a Measured Wave Field", "The route deployable today")
bullets(s, [
    {"t": "A receiver measures the channel \u2014 the environment\u2019s transformation of the signal:", "b": False},
], Inches(0.85), Inches(1.75), Inches(11.6), Inches(0.7), size=17)
picture(s, "eq_csi_dyn.png", Inches(3.0), Inches(2.45), w=Inches(7.2))
bullets(s, [
    {"t": "The static term carries fixed geometry and motionless people; the moving term carries Doppler from motion.", "b": False},
    {"t": "The CSI tensor (antenna \u00d7 subcarrier \u00d7 time) is the wireless \u201cpixel grid\u201d \u2014 a direct sample of the wave field.", "b": True},
    {"t": "A network learns the inverse map: through-wall pose, body surfaces from Wi-Fi, gestures, activity (published systems).", "b": False},
], Inches(0.85), Inches(3.85), Inches(11.6), Inches(2.6), size=16.5, gap=11)
footer(s, 8, "From Light to Radio  \u00b7  Channel & CSI model")

# ===========================================================================
# 9. META-PERCEPTION
# ===========================================================================
s = slide()
title(s, "Meta-Perception: Vision + Wireless", "Complementary failure modes")
bullets(s, [
    {"t": "Vision: high detail, but fails in dark / glare / occlusion; privacy-heavy; costly always-on.", "b": False},
    {"t": "Wireless: coarse, but cheap, private, penetrating, always-on.", "b": False},
    {"t": "Fuse as a posterior over the scene; the confident modality dominates automatically.", "b": True},
], Inches(0.85), Inches(1.75), Inches(5.5), Inches(3.6), size=16.5, gap=12)
picture(s, "eq_fusion.png", Inches(0.85), Inches(5.25), w=Inches(5.3))
picture(s, "fig_fusion.png", Inches(6.5), Inches(2.1), w=Inches(6.4))
caption(s, "Wireless runs always-on; vision wakes on confident events.", Inches(6.5), Inches(5.45), Inches(6.4))
footer(s, 9, "From Light to Radio  \u00b7  Fusion architecture")

# ===========================================================================
# 10. APPLICATIONS TABLE
# ===========================================================================
s = slide()
title(s, "Applications and Directions", "Where this goes \u2014 and how to choose one")
rows_data = [
    ["Direction", "Example", "Data required", "Why fuse"],
    ["Elderly / home care", "Fall detection, 24/7 presence", "Wi-Fi CSI + sparse camera labels", "Private, works in the dark"],
    ["Security / rescue", "Through-wall in smoke, dark", "mmWave / UWB radar + IR", "RF penetrates; vision adds ID"],
    ["Automotive / robotics", "Perception in fog, glare", "Radar point cloud + camera", "Radar robust; vision semantics"],
    ["Vital signs", "Contactless breathing / heart-rate", "Fine-phase CSI / FMCW radar", "Continuous; vision for context"],
    ["Smart buildings", "Occupancy for HVAC, lighting", "Coarse CSI presence per zone", "Cheap counting; camera audits"],
]
nrows, ncols = len(rows_data), 4
gtbl = s.shapes.add_table(nrows, ncols, Inches(0.7), Inches(1.7), Inches(11.95), Inches(4.6)).table
gtbl.columns[0].width = Inches(2.5)
gtbl.columns[1].width = Inches(3.2)
gtbl.columns[2].width = Inches(3.35)
gtbl.columns[3].width = Inches(2.9)
for ri, row in enumerate(rows_data):
    gtbl.rows[ri].height = Inches(0.5)
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = HEADER_FILL
        else:
            cell.fill.fore_color.rgb = ROW_FILL if ri % 2 == 0 else WHITE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val
        _set_font(r, 11.5 if ri == 0 else 11, INK if ri else SLATE, bold=(ri == 0 or ci == 0))
bullets(s, [
    {"t": "Down-select on data availability, resolution sufficiency, privacy fit and cost \u2014 care and occupancy usually win.", "b": True},
], Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.6), size=13, gap=4)
footer(s, 10, "From Light to Radio  \u00b7  Application directions")

# ===========================================================================
# 11. DATA REQUIREMENTS
# ===========================================================================
s = slide()
title(s, "Sizing the Data, From First Principles", "Four back-of-envelope formulas")
bullets(s, [
    {"t": "Spatial detail (the wireless pixel count):  N_dof \u2248 A_T A_R / (\u03bb d)\u00b2", "b": True},
    {"t": "Range from bandwidth:  \u0394R = c / 2B   \u2014  30 cm needs \u2248 0.5 GHz", "b": False},
    {"t": "Angle from aperture:  \u0394\u03b8 \u2248 \u03bb / D ;  velocity from time:  \u0394v \u2248 \u03bb / 2T", "b": False},
    {"t": "How many measurements (sparse scene):  M \u2265 C\u00b7S\u00b7log(N/S)", "b": True},
    {"t": "Training: a few thousand labelled windows, auto-labelled by a camera during a short enrolment.", "b": False},
], Inches(0.85), Inches(1.75), Inches(6.0), Inches(4.7), size=16, gap=13)
picture(s, "fig_data_requirements.png", Inches(7.05), Inches(2.25), w=Inches(5.9))
caption(s, "Read bandwidth off the range target; read spatial detail off aperture.", Inches(7.05), Inches(4.95), Inches(5.9))
footer(s, 11, "From Light to Radio  \u00b7  Data requirements")

# ===========================================================================
# 12. EVENT CAMERAS VS ALWAYS-ON WIRELESS
# ===========================================================================
s = slide()
title(s, "Event Cameras vs Always-On Wireless", "Why wireless is right for constant monitoring")
bullets(s, [
    {"t": "Event cameras fire only on change \u2014 silent on a still person (asleep, seated, collapsed).", "b": False},
    {"t": "Wireless is active: the transmitter always illuminates the scene.", "b": True},
    {"t": "A motionless person still shifts the static channel \u2014 detectable vs an empty-room reference.", "b": False},
    {"t": "We differ by staying always-on: presence and motion at once; change-detection only gates vision.", "b": True},
], Inches(0.85), Inches(1.75), Inches(5.6), Inches(4.6), size=16, gap=12)
picture(s, "fig_event_vs_wireless.png", Inches(6.7), Inches(1.95), w=Inches(6.2))
caption(s, "Active wireless sees a still person and motion alike.", Inches(6.7), Inches(5.55), Inches(6.2))
footer(s, 12, "From Light to Radio  \u00b7  Continuous monitoring")

# ===========================================================================
# 13. CONCLUSION
# ===========================================================================
s = slide()
title(s, "Bottom Line", "Is it possible, mathematically? Yes.")
bullets(s, [
    {"t": "Light and radio are the same wave \u2014 same Helmholtz equation, different wavelength.", "b": True},
    {"t": "Scale invariance: the maths of optical wave-perception transfers to wireless exactly.", "b": False},
    {"t": "The one trade-off is resolution \u2014 coarse, but private and penetrating.", "b": False},
    {"t": "The reward: continuous, privacy-preserving, all-conditions perception that complements vision.", "b": True},
], Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.0), size=19, gap=16)
footer(s, 13, "From Light to Radio  \u00b7  Conclusion")

# ===========================================================================
# 14. REFERENCES (selected)
# ===========================================================================
s = slide()
title(s, "Selected References", "Peer-reviewed, with DOIs")
refs = [
    "Lin et al. (2018). All-optical machine learning using diffractive deep neural networks. Science 361, 1004\u20131008.",
    "Chen, Li, Wang, Chen & Ozcan (2025). Optical generative models. Nature 644, 903\u2013911.",
    "Li et al. (2019). Machine-learning reprogrammable metasurface imager. Nature Communications 10, 1082.",
    "Li et al. (2019). Intelligent metasurface imager and recognizer. Light: Science & Applications 8, 97.",
    "Hunt et al. (2013). Metamaterial apertures for computational imaging. Science 339, 310\u2013313.",
    "Zhao et al. (2018). Through-wall human pose estimation using radio signals. CVPR 2018.",
    "Geng, Huang & De la Torre (2023). DensePose from WiFi. arXiv:2301.00250.",
    "Yang et al. (2023). SenseFi: deep-learning-empowered WiFi human sensing. Patterns 4, 100703.",
    "Ma, Zhou & Wang (2019). WiFi sensing with channel state information: a survey. ACM CSUR 52(3).",
    "Cand\u00e8s & Wakin (2008). An introduction to compressive sampling. IEEE SPM 25(2), 21\u201330.",
    "Gallego et al. (2022). Event-based vision: a survey. IEEE TPAMI 44(1), 154\u2013180.",
]
tb, tf = textbox(s, Inches(0.85), Inches(1.8), Inches(11.8), Inches(5.0))
for i, ref in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(7); p.line_spacing = 1.02
    rn = p.add_run(); rn.text = f"[{i+1}]  "
    _set_font(rn, 12, SLATE, bold=True)
    r = p.add_run(); r.text = ref
    _set_font(r, 12, INK)
footer(s, 14, "From Light to Radio  \u00b7  Full reference list in the white paper")

prs.save(OUT)
print("SAVED", OUT)
